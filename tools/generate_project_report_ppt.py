# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=1.0.2",
# ]
# ///

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).resolve().parent.parent / "项目汇报_序时账审计分析平台_浅色商务版.pptx"

BG = RGBColor(247, 248, 250)
PANEL = RGBColor(255, 255, 255)
PANEL_ALT = RGBColor(238, 242, 247)
NAVY = RGBColor(29, 57, 93)
TEAL = RGBColor(88, 127, 140)
ACCENT = RGBColor(197, 131, 74)
TEXT = RGBColor(45, 55, 72)
MUTED = RGBColor(104, 118, 138)
LINE = RGBColor(217, 224, 233)
SOFT_GREEN = RGBColor(225, 237, 233)
SOFT_BLUE = RGBColor(229, 236, 245)
SOFT_ORANGE = RGBColor(248, 237, 226)

FONT = "PingFang SC"
FONT_EN = "Aptos"


def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_textbox(slide, x, y, w, h, text, *, size=16, color=TEXT, bold=False,
                font=FONT, align=PP_ALIGN.LEFT, margin=0.08, valign=MSO_ANCHOR.TOP,
                line_spacing=1.18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(slide, x, y, w, h, lines, *, size=15, color=TEXT, font=FONT,
                  bullet=False, gap=2, margin=0.06, first_bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {line}" if bullet else line
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.18
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = first_bold and idx == 0
            run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, *, fill=PANEL, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_section_chip(slide, x, y, text, *, fill=SOFT_BLUE, color=NAVY, w=1.18):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.color.rgb = fill
    add_textbox(slide, x + 0.02, y + 0.02, w - 0.04, 0.28, text, size=11, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_title_block(slide, page_no, section, title, subtitle):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.28))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = NAVY
    top_band.line.color.rgb = NAVY
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.58), Inches(0.12), Inches(0.82))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.color.rgb = ACCENT
    add_section_chip(slide, 8.55, 0.6, section)
    add_textbox(slide, 0.92, 0.5, 6.9, 0.55, title, size=28, color=NAVY, bold=True, font=FONT)
    if subtitle:
        add_textbox(slide, 0.94, 1.02, 7.8, 0.38, subtitle, size=11.5, color=MUTED, font=FONT)
    add_textbox(slide, 9.15, 5.12, 0.45, 0.22, f"{page_no:02d}", size=10, color=MUTED, font=FONT_EN, align=PP_ALIGN.RIGHT)


def add_cover_title(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.62), Inches(8.76), Inches(4.15))
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.color.rgb = LINE
    side_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.62), Inches(0.2), Inches(4.15))
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = NAVY
    side_bar.line.color.rgb = NAVY
    add_textbox(slide, 1.06, 0.95, 6.8, 0.55, "序时账审计分析平台", size=30, color=NAVY, bold=True)
    add_textbox(slide, 1.08, 1.52, 7.1, 0.6, "项目汇报 PPT", size=20, color=ACCENT, bold=True, font=FONT_EN)
    add_textbox(
        slide,
        1.08,
        2.0,
        6.8,
        0.95,
        "面向多年序时账的辅助分析工具，围绕可视化、规则抽样和 LLM 辅助核实，帮助审计人员更快完成前期数据初筛与样本定位。",
        size=16,
        color=TEXT,
    )
    add_textbox(
        slide,
        1.08,
        4.08,
        6.2,
        0.34,
        "风格定位：浅色商务风 | 叙述口径：审慎、可解释、不夸大",
        size=11.5,
        color=MUTED,
    )
    steps = ["上传数据", "数据画像", "规则校准", "执行抽样", "报告输出"]
    for idx, step in enumerate(steps):
        x = 6.95 + (idx % 2) * 1.1
        y = 1.18 + (idx // 2) * 1.0
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL if idx < 4 else ACCENT
        circ.line.color.rgb = circ.fill.fore_color.rgb
        add_textbox(slide, x, y + 0.02, 0.45, 0.3, str(idx + 1), size=12, color=PANEL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        add_card(slide, x + 0.53, y - 0.02, 1.35, 0.5, fill=PANEL_ALT, line=LINE)
        add_textbox(slide, x + 0.62, y + 0.08, 1.15, 0.22, step, size=11, color=TEXT, bold=True, margin=0)


def add_bullet_card(slide, x, y, w, h, title, lines, *, fill=PANEL, accent=NAVY):
    add_card(slide, x, y, w, h, fill=fill)
    accent_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.color.rgb = accent
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.3, 0.28, title, size=16, color=NAVY, bold=True)
    add_multiline(slide, x + 0.18, y + 0.54, w - 0.28, h - 0.66, lines, size=13.2, color=TEXT, bullet=True)


def add_metric_card(slide, x, y, w, h, value, label, note, *, fill=PANEL):
    add_card(slide, x, y, w, h, fill=fill)
    add_textbox(slide, x + 0.12, y + 0.12, w - 0.24, 0.34, value, size=22, color=NAVY, bold=True, font=FONT_EN)
    add_textbox(slide, x + 0.12, y + 0.52, w - 0.24, 0.24, label, size=11.8, color=TEXT, bold=True)
    add_textbox(slide, x + 0.12, y + 0.8, w - 0.24, h - 0.92, note, size=10.5, color=MUTED)


def add_flow_step(slide, x, y, w, h, step_no, title, body, *, fill=PANEL):
    add_card(slide, x, y, w, h, fill=fill)
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.color.rgb = TEAL
    add_textbox(slide, x + 0.18, y + 0.2, 0.42, 0.26, str(step_no), size=11.5, color=PANEL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_textbox(slide, x + 0.7, y + 0.16, w - 0.88, 0.28, title, size=14, color=NAVY, bold=True)
    add_textbox(slide, x + 0.18, y + 0.68, w - 0.3, h - 0.78, body, size=11.5, color=TEXT)


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = LINE
    line.line.width = Pt(2)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    prs.core_properties.author = "Codex"
    prs.core_properties.title = "序时账审计分析平台项目汇报"
    prs.core_properties.subject = "浅色商务风项目汇报"

    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_cover_title(slide)

    # 2 Project positioning
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 2, "定位", "项目定位与目标", "先说明项目是做什么的，再说明它解决到什么程度。")
    add_bullet_card(slide, 0.7, 1.55, 2.7, 2.75, "项目定位", [
        "面向多年序时账数据的辅助分析平台。",
        "将上传、画像、规则抽样和报告输出串成完整流程。",
        "更适合作为审计前期分析和样本定位工具使用。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.6, 1.55, 2.7, 2.75, "项目目的", [
        "先建立对账套规模、结构和波动的整体理解。",
        "再把值得继续穿透的位置收缩到有限范围。",
        "减少人工直接逐表浏览的工作量。",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 6.5, 1.55, 2.8, 2.75, "作用边界", [
        "不替代审计判断，也不自动形成结论。",
        "不覆盖全部业务背景和全部证据链。",
        "核心作用是提高前期分析与抽样效率。",
    ], fill=SOFT_ORANGE, accent=ACCENT)

    # 3 Why visualization
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 3, "可视化", "为什么项目先做可视化", "这一页强调：可视化的任务是帮助形成整体判断，不是直接给结论。")
    add_metric_card(slide, 0.72, 1.55, 2.15, 1.4, "先看模式", "不是先看凭证", "序时账原始明细通常行数较大，直接逐行阅读很难快速把握重点。", fill=SOFT_BLUE)
    add_metric_card(slide, 3.05, 1.55, 2.15, 1.4, "先看波动", "不是先下结论", "先识别月份、科目、对手方和操作人层面的异常位置。", fill=SOFT_GREEN)
    add_metric_card(slide, 5.38, 1.55, 2.15, 1.4, "先缩范围", "不是全部展开", "通过图表把后续需要穿透的范围缩小到相对可处理的集合。", fill=SOFT_ORANGE)
    add_metric_card(slide, 7.71, 1.55, 1.6, 1.4, "辅助", "不是替代", "帮助做初筛与排序。", fill=PANEL_ALT)
    add_bullet_card(slide, 0.72, 3.2, 8.6, 1.45, "本项目中可视化的真实作用", [
        "帮助回答三个问题：数据整体长什么样、哪些位置波动较大、哪些位置值得继续下钻。",
        "为后续规则校准、抽样执行和人工复核提供可解释的前置信息。",
        "图表本身不直接证明异常成立，后续仍需凭证、合同、审批和期后证据支持。",
    ], fill=PANEL, accent=NAVY)

    # 4 Workflow
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 4, "流程", "项目整体流程", "围绕 5 步工作流展开，平台当前实现与 Streamlit 主流程一致。")
    xs = [0.68, 2.62, 4.56, 6.5, 8.04]
    titles = ["上传数据", "数据画像", "规则校准", "执行抽样", "报告输出"]
    bodies = [
        "读取单年、多年或多文件 Excel，统一字段并识别年度。",
        "生成财务概况、统计画像、审计可视化和跨年发现。",
        "结合画像摘要、跨年线索和经验库生成项目级规则配置。",
        "按规则扫描分录，形成候选样本池，并可选做 LLM 辅助核实。",
        "输出样本清单、规则统计和命中明细，供人工继续处理。",
    ]
    for i, (x, title, body) in enumerate(zip(xs, titles, bodies, strict=True), start=1):
        add_flow_step(slide, x, 1.75, 1.45, 2.55, i, title, body, fill=PANEL if i % 2 else PANEL_ALT)
        if i < 5:
            add_connector(slide, x + 1.45, 3.02, xs[i] - 0.08, 3.02)

    # 5 Visualization overview
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 5, "模块", "可视化模块总览", "图表不是堆在一起展示，而是按审计分析问题拆成几个相对清晰的入口。")
    modules = [
        ("财务概况", "先看企业经营轮廓和年度结构。", SOFT_BLUE),
        ("收入成本", "看月度波动和借贷方向异常。", SOFT_GREEN),
        ("费用分析", "看费用结构和异常放量类别。", SOFT_ORANGE),
        ("客户供应商", "看交易对手方集中度和重点对象。", PANEL_ALT),
        ("暂估与往来", "看挂账形成、核销和净额变化。", SOFT_BLUE),
        ("调账冲销", "集中查看调账、冲销、反记账凭证。", SOFT_GREEN),
    ]
    idx = 0
    for row in range(2):
        for col in range(3):
            x = 0.72 + col * 3.02
            y = 1.62 + row * 1.55
            title, body, fill = modules[idx]
            add_card(slide, x, y, 2.72, 1.2, fill=fill)
            add_textbox(slide, x + 0.14, y + 0.16, 2.4, 0.26, title, size=15, color=NAVY, bold=True)
            add_textbox(slide, x + 0.14, y + 0.5, 2.42, 0.5, body, size=11.2, color=TEXT)
            idx += 1
    add_bullet_card(slide, 0.72, 4.88 - 0.62, 8.6, 0.78, "设计原因", [
        "按问题拆分模块，方便先进入“要解决什么分析任务”，而不是先面对一整页明细和图表。"
    ], fill=PANEL, accent=ACCENT)

    # 6 Financial overview
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 6, "可视化", "财务概况模块：先建立经营背景", "对应财务概况 Tab，先帮助使用者形成对企业收入、成本和费用结构的基本认识。")
    add_bullet_card(slide, 0.72, 1.55, 3.0, 2.8, "为什么设置这个模块", [
        "后续看到异常时，需要先知道企业的大致经营规模和主要构成。",
        "如果财务轮廓本身不清晰，后面看到波动也缺少解释背景。",
        "因此它更多提供理解上下文，而不是直接给出风险结论。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.94, 1.55, 2.46, 2.8, "当前展示内容", [
        "总收入、总成本、毛利率、研发费用率。",
        "多年财务对比摘要。",
        "成本结构、收入构成、投资收益与营业外收入。",
    ], fill=PANEL, accent=TEAL)
    add_card(slide, 6.62, 1.55, 2.68, 2.8, fill=SOFT_GREEN)
    add_textbox(slide, 6.78, 1.72, 2.1, 0.24, "在汇报里的表述建议", size=14, color=NAVY, bold=True)
    add_textbox(slide, 6.78, 2.08, 2.18, 1.72, "这一模块的作用，是先建立对企业年度经营轮廓的整体认识，帮助后续理解哪些月份、科目和对象更值得继续穿透。", size=12.6, color=TEXT)

    # 7 Revenue cost
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 7, "可视化", "收入成本模块：看波动，也看方向", "平台不仅看净额，还保留借贷方向拆分，避免异常被净额掩盖。")
    add_card(slide, 0.72, 1.55, 4.02, 2.7, fill=PANEL)
    add_textbox(slide, 0.88, 1.72, 2.6, 0.24, "为什么设置这个模块", size=15, color=NAVY, bold=True)
    add_multiline(slide, 0.88, 2.08, 3.55, 1.82, [
        "收入和成本是利润形成的核心，通常也是审计优先关注区域。",
        "仅看净收入和净成本不够，部分异常会体现在收入 S、成本 H 这类方向性信号上。",
        "因此该模块更适合先定位“哪里可能有问题”，再决定是否下钻到分录。",
    ], bullet=True, size=12.6)
    add_card(slide, 5.0, 1.55, 4.3, 2.7, fill=SOFT_ORANGE)
    add_textbox(slide, 5.16, 1.72, 2.4, 0.24, "当前展示内容", size=15, color=NAVY, bold=True)
    add_multiline(slide, 5.16, 2.08, 3.8, 1.68, [
        "月度净收入、净成本、毛利趋势。",
        "收入 S / 成本 H 异常方向对比。",
        "点击月份后查看相关 Top10 分录。",
        "总览页同时给出净收入、净成本、毛利和异常方向总额。",
    ], bullet=True, size=12.6)
    add_bullet_card(slide, 0.72, 4.42, 8.58, 0.86, "汇报建议", [
        "这一模块主要服务于发现月度层面的波动和方向性异常，不直接说明原因，后续仍需结合具体分录和业务背景解释。"
    ], fill=PANEL_ALT, accent=ACCENT)

    # 8 Expense
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 8, "可视化", "费用模块：从总额切到类别", "费用分析不是只看费用有多少，而是看哪些费用类型更值得做实质性测试。")
    add_metric_card(slide, 0.72, 1.58, 2.0, 1.42, "结构", "费用构成", "先看费用由哪些主要类别组成。", fill=SOFT_BLUE)
    add_metric_card(slide, 2.96, 1.58, 2.0, 1.42, "占比", "费用轻重", "看哪些费用类别占比更高。", fill=SOFT_GREEN)
    add_metric_card(slide, 5.2, 1.58, 2.0, 1.42, "放量", "异常类别", "看哪些费用类别同比或结构上更值得关注。", fill=SOFT_ORANGE)
    add_metric_card(slide, 7.44, 1.58, 1.86, 1.42, "下钻", "分录查看", "点击类别后直接查看分录。", fill=PANEL_ALT)
    add_bullet_card(slide, 0.72, 3.28, 4.08, 1.78, "为什么设置这个模块", [
        "费用总额本身信息量有限，需要拆到类别层级才更方便判断。",
        "费用类项目常涉及利润调节、不合规支出或利益输送等后续核查方向。",
        "平台当前只负责先暴露类别差异和重点样本入口。",
    ], fill=PANEL, accent=NAVY)
    add_bullet_card(slide, 5.02, 3.28, 4.28, 1.78, "当前展示内容", [
        "费用结构图、费用类别金额与占比表。",
        "研发费用、财务费用、税金及附加等补充项。",
        "点击费用类别后查看 Top10 分录。",
    ], fill=PANEL, accent=TEAL)

    # 9 Customer supplier
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 9, "可视化", "客户供应商模块：先锁定重点对象", "很多后续核查最终都会落到交易对手方，因此平台单独保留这一视角。")
    add_card(slide, 0.72, 1.55, 4.1, 2.75, fill=SOFT_BLUE)
    add_textbox(slide, 0.88, 1.72, 1.9, 0.24, "客户收入 Top10", size=15, color=NAVY, bold=True)
    add_multiline(slide, 0.88, 2.08, 3.6, 1.65, [
        "帮助先识别收入集中客户。",
        "便于判断是否存在收入依赖、关联交易或年底冲量入口。",
        "点击后可直接查看对应分录 Top10。",
    ], bullet=True, size=12.4)
    add_card(slide, 5.12, 1.55, 4.18, 2.75, fill=SOFT_GREEN)
    add_textbox(slide, 5.28, 1.72, 2.1, 0.24, "供应商应付 Top10", size=15, color=NAVY, bold=True)
    add_multiline(slide, 5.28, 2.08, 3.65, 1.65, [
        "帮助先识别采购或往来集中供应商。",
        "便于判断哪些对象更值得继续看付款与挂账明细。",
        "展示结果是对象优先级，不是风险定性。",
    ], bullet=True, size=12.4)
    add_bullet_card(slide, 0.72, 4.52, 8.58, 0.74, "设计原因", [
        "客户和供应商维度的主要价值，在于把“后面优先看谁”尽快明确下来。"
    ], fill=PANEL_ALT, accent=ACCENT)

    # 10 Working capital
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 10, "可视化", "暂估与往来模块：观察形成、核销与净额变化", "往来类科目单看年末余额不够，因此平台按月保留借贷方向和净额轨迹。")
    add_bullet_card(slide, 0.72, 1.55, 2.72, 2.82, "应付账款暂估", [
        "看暂估贷方增加、借方减少和净额变化。",
        "更适合观察暂估是如何形成、如何消化。",
        "点击月份可继续看供应商和分录。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.64, 1.55, 2.72, 2.82, "其他应收", [
        "保留 S、H 和净额，不直接压成单一余额。",
        "帮助观察是否存在持续挂账或异常波动月份。",
        "适合作为后续往来核查的线索入口。",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 6.56, 1.55, 2.74, 2.82, "其他应付", [
        "保留预提 H、核销 S 和净值变化。",
        "便于区分形成和冲回，而不是只看一个时点数字。",
        "更适合辅助识别继续检查的月份。",
    ], fill=SOFT_ORANGE, accent=ACCENT)
    add_bullet_card(slide, 0.72, 4.52, 8.58, 0.74, "设计原因", [
        "这一模块的重点不在“证明异常”，而在于把往来类科目的时间轨迹先展示出来。"
    ], fill=PANEL, accent=NAVY)

    # 11 Adjustment reversal
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 11, "可视化", "调账冲销模块：集中查看高敏感凭证", "调账、冲销、反记账类凭证未必一定异常，但通常更值得优先复核。")
    add_card(slide, 0.72, 1.58, 2.42, 2.7, fill=SOFT_ORANGE)
    add_textbox(slide, 0.88, 1.78, 1.9, 0.24, "关键词入口", size=15, color=NAVY, bold=True)
    add_multiline(slide, 0.88, 2.14, 2.0, 1.55, [
        "冲账", "冲销", "冲回", "调账", "反记账", "红字", "重分类", "更正"
    ], bullet=True, size=12.4)
    add_bullet_card(slide, 3.36, 1.58, 2.68, 2.7, "为什么设置这个模块", [
        "这类凭证通常具有较强审计关注属性。",
        "单独集中展示，便于统一查看和下钻。",
        "目的是优先排序，不是直接认定问题。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 6.28, 1.58, 3.02, 2.7, "当前展示内容", [
        "调账冲销凭证汇总表。",
        "按关键词筛选相关凭证。",
        "查看同凭证完整分录。",
    ], fill=SOFT_GREEN, accent=TEAL)

    # 12 Statistical profile and cross-year
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 12, "画像", "统计画像与跨年交叉稽核", "这一层更多服务于建立账套基线和补充单年观察的局限。")
    add_bullet_card(slide, 0.72, 1.55, 4.14, 2.75, "统计画像的作用", [
        "建立账套规模、凭证类型、用户集中度和金额分布基线。",
        "当前展示包括总体规模、金额分位数、本福特分布、月末集中度等。",
        "主要作为后续规则校准的输入之一。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 5.1, 1.55, 4.2, 2.75, "跨年稽核的作用", [
        "补足单年视角下不容易识别的问题。",
        "当前覆盖预提冲回、收入跨年、余额累积、资金循环等类型。",
        "输出的是跨年线索，不替代截止测试或往来核查程序。",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 0.72, 4.5, 8.58, 0.76, "汇报建议", [
        "可以把这一页理解成“平台的基线层和趋势层”，它们主要负责给后面的规则和样本提供背景。"
    ], fill=PANEL_ALT, accent=ACCENT)

    # 13 Architecture
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 13, "架构", "项目主要组成部分", "平台内部按数据摄入、分析、规则、核实和输出分模块实现。")
    arch = [
        ("数据摄入", "ingestion.py", "读取 Excel、统一字段、识别年度、标记 Period 13。"),
        ("画像分析", "profiler.py / visual_analysis.py", "生成财务概况、统计画像和审计可视化。"),
        ("跨年分析", "cross_year.py", "执行跨年异常检测并形成发现。"),
        ("规则校准", "rule_generator.py / knowledge_base.py", "汇总摘要并生成项目级 rules_config。"),
        ("抽样执行", "rule_engine.py", "执行规则扫描并形成命中结果。"),
        ("辅助核实", "llm_verifier.py / reporter.py", "做可选 LLM 核实并输出 Excel 报告。"),
    ]
    for i, (name, mod, body) in enumerate(arch):
        col = i % 2
        row = i // 2
        x = 0.72 + col * 4.4
        y = 1.55 + row * 1.08
        add_card(slide, x, y, 4.0, 0.9, fill=PANEL if i % 2 == 0 else PANEL_ALT)
        add_textbox(slide, x + 0.14, y + 0.12, 1.35, 0.2, name, size=14, color=NAVY, bold=True)
        add_textbox(slide, x + 1.56, y + 0.12, 2.18, 0.2, mod, size=10.4, color=ACCENT, bold=True, font=FONT_EN)
        add_textbox(slide, x + 0.14, y + 0.42, 3.62, 0.28, body, size=11.1, color=TEXT)

    # 14 Rule calibration overview
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 14, "规则", "规则校准：目标与输入", "规则校准的任务，是把通用规则骨架调成更贴近当前账套特征的项目级配置。")
    add_bullet_card(slide, 0.72, 1.55, 2.75, 2.75, "校准目标", [
        "调整规则启用状态和阈值。",
        "让抽样参数更贴近当前账套环境。",
        "避免完全依赖固定经验阈值。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.66, 1.55, 2.75, 2.75, "输入一：画像摘要", [
        "财务概况摘要",
        "统计画像摘要",
        "金额分布、用户集中度、手工凭证占比等",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 6.6, 1.55, 2.7, 2.75, "输入二、三", [
        "跨年异常摘要",
        "经验库历史有效规则",
        "历史确认率、参数与说明",
    ], fill=SOFT_ORANGE, accent=ACCENT)

    # 15 Calibration output and boundaries
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 15, "规则", "规则校准：输出结构与作用边界", "平台输出的是结构化规则配置建议，后续仍保留人工查看和修改空间。")
    add_card(slide, 0.72, 1.55, 3.34, 2.78, fill=PANEL)
    add_textbox(slide, 0.88, 1.72, 2.2, 0.22, "输出内容", size=15, color=NAVY, bold=True)
    add_multiline(slide, 0.88, 2.08, 2.9, 1.72, [
        "是否启用 enabled",
        "阈值参数",
        "rationale 说明",
        "白名单与样本量控制",
    ], bullet=True, size=12.6)
    add_card(slide, 4.28, 1.55, 5.02, 2.78, fill=SOFT_BLUE)
    add_textbox(slide, 4.46, 1.72, 2.4, 0.22, "作用边界", size=15, color=NAVY, bold=True)
    add_multiline(slide, 4.46, 2.08, 4.46, 1.7, [
        "模型看到的是摘要而不是全部底层证据。",
        "当前更适合作为参数建议生成器，而不是无条件自动定规则。",
        "当账套存在特殊业务模式时，仍然需要人工判断是否关闭、放宽或收紧某些规则。",
    ], bullet=True, size=12.4)
    add_bullet_card(slide, 0.72, 4.52, 8.58, 0.74, "汇报建议", [
        "建议把“规则校准”表述为项目级参数调整过程，而不是说成模型已经自动理解业务并生成正确规则。"
    ], fill=PANEL_ALT, accent=ACCENT)

    # 16 Rule scope
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 16, "规则", "当前纳入校准与执行的规则范围", "规则范围代表当前版本已经结构化实现的部分，不代表覆盖全部审计风险。")
    rules = [
        "化整为零", "大额异常", "手工凭证", "计提异常",
        "收入突增", "融资性贸易", "资金池划转", "用户集中度异常",
        "冲销反记账异常", "敏感费用筛查", "跨年异常相关规则", "白名单与样本上限"
    ]
    for i, rule in enumerate(rules):
        row = i // 4
        col = i % 4
        x = 0.72 + col * 2.18
        y = 1.62 + row * 0.92
        fill = [SOFT_BLUE, SOFT_GREEN, SOFT_ORANGE, PANEL_ALT][col]
        add_card(slide, x, y, 1.92, 0.66, fill=fill)
        add_textbox(slide, x + 0.1, y + 0.19, 1.72, 0.2, rule, size=11.3, color=TEXT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)

    # 17 Sampling logic
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 17, "抽样", "抽样执行逻辑", "规则执行阶段先做过滤，再生成命中项，最后汇总成候选样本池。")
    steps = [
        ("01 过滤", "应用白名单关键词、白名单凭证类型，并过滤金额为 0 的记录。"),
        ("02 执行", "各规则函数独立扫描统一后的 DataFrame。"),
        ("03 命中", "每个命中转成 RuleHit，记录凭证号、规则类型、证据和优先级。"),
        ("04 汇总", "跨年发现也会转成命中项，一起进入样本池。"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.72 + i * 2.25
        add_flow_step(slide, x, 1.9, 1.9, 2.2, i + 1, title, body, fill=PANEL if i % 2 == 0 else PANEL_ALT)
        if i < 3:
            add_connector(slide, x + 1.9, 3.0, x + 2.16, 3.0)

    # 18 Sorting and report
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 18, "抽样", "样本排序、截断与报告输出", "平台不会直接把所有命中结果原样抛出，而是保留样本量上限和排序逻辑。")
    add_bullet_card(slide, 0.72, 1.55, 2.74, 2.78, "排序逻辑", [
        "无 LLM 时，主要依据规则优先级排序。",
        "有 LLM 时，可依据核实后的风险级别优先输出。",
        "同一凭证会先去重，再进入样本清单。",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.66, 1.55, 2.74, 2.78, "截断逻辑", [
        "保留样本量上限配置。",
        "避免输出过多凭证造成二次处理负担。",
        "更适合人工继续复核和审计程序安排。",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 6.6, 1.55, 2.7, 2.78, "输出内容", [
        "样本清单",
        "规则统计",
        "命中明细",
        "可选 LLM 核实字段",
    ], fill=SOFT_ORANGE, accent=ACCENT)

    # 19 LLM verification flow
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 19, "LLM", "大模型辅助核实：流程与定位", "这一阶段的目标是辅助排序和补充核查建议，不是自动给出最终审计结论。")
    add_metric_card(slide, 0.72, 1.55, 2.0, 1.48, "Top 50", "候选凭证", "当前配置下，对优先级最高的有限数量凭证做辅助核实。", fill=SOFT_BLUE)
    add_metric_card(slide, 2.96, 1.55, 2.0, 1.48, "按凭证", "先去重", "避免同一凭证因命中多条规则被重复送入模型。", fill=SOFT_GREEN)
    add_metric_card(slide, 5.2, 1.55, 2.0, 1.48, "分批", "批量核实", "当前默认按批次送入，降低单次输入规模。", fill=SOFT_ORANGE)
    add_metric_card(slide, 7.44, 1.55, 1.86, 1.48, "可选", "不是必开", "未启用时平台仍可完成抽样与报告输出。", fill=PANEL_ALT)
    add_bullet_card(slide, 0.72, 3.35, 8.58, 1.8, "定位说明", [
        "模型看到的是规则命中证据和部分分录摘要，信息完整性有限。",
        "因此这一步更适合作为二次筛选器，帮助区分哪些凭证更值得优先复核。",
        "汇报中建议明确写成“辅助核实”或“辅助判断”，避免说成自动审计或自动定性。",
    ], fill=PANEL, accent=NAVY)

    # 20 LLM input output boundaries
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 20, "LLM", "辅助核实的输入、输出与保守策略", "当前实现强调保守处理：即使模型失败，也尽量不丢掉已命中的候选样本。")
    add_bullet_card(slide, 0.72, 1.55, 2.72, 2.78, "输入内容", [
        "凭证编号",
        "规则类型",
        "触发证据",
        "部分分录行信息",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.64, 1.55, 2.72, 2.78, "输出内容", [
        "是否确认",
        "风险级别",
        "判断理由",
        "建议核查程序",
    ], fill=SOFT_GREEN, accent=TEAL)
    add_bullet_card(slide, 6.56, 1.55, 2.74, 2.78, "保守策略", [
        "失败时不直接删除样本。",
        "保留为待复核状态。",
        "优先保证样本不丢失。",
    ], fill=SOFT_ORANGE, accent=ACCENT)

    # 21 Knowledge base and close loop
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 21, "闭环", "经验库与项目闭环", "平台输出不仅面向本次项目，也保留了一部分跨项目复用能力。")
    add_flow_step(slide, 0.72, 1.72, 2.0, 2.45, 1, "本次项目规则", "先在当前账套中完成参数校准与样本执行。", fill=SOFT_BLUE)
    add_flow_step(slide, 3.06, 1.72, 2.0, 2.45, 2, "人工复核结果", "由审计人员判断哪些规则在本项目中更有效。", fill=SOFT_GREEN)
    add_flow_step(slide, 5.4, 1.72, 2.0, 2.45, 3, "经验库沉淀", "把有效规则的参数、说明和确认表现保存起来。", fill=SOFT_ORANGE)
    add_flow_step(slide, 7.74, 1.72, 1.56, 2.45, 4, "后续项目参考", "供下一次规则校准时作为历史输入。", fill=PANEL_ALT)
    add_connector(slide, 2.72, 2.95, 3.02, 2.95)
    add_connector(slide, 5.06, 2.95, 5.36, 2.95)
    add_connector(slide, 7.4, 2.95, 7.7, 2.95)

    # 22 Boundaries
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 22, "边界", "项目当前定位与适用边界", "这一页建议在汇报末尾明确说清楚，整体口径会更稳。")
    add_bullet_card(slide, 0.72, 1.55, 2.82, 2.86, "当前更适合", [
        "审计前期数据初筛",
        "项目级规则参数调整",
        "重点样本定位与排序",
        "结果导出和人工继续复核",
    ], fill=SOFT_BLUE, accent=NAVY)
    add_bullet_card(slide, 3.74, 1.55, 2.82, 2.86, "当前不宜夸大为", [
        "自动审计系统",
        "自动定性结论系统",
        "覆盖全部业务风险的规则库",
        "替代人工程序的核实工具",
    ], fill=SOFT_ORANGE, accent=ACCENT)
    add_bullet_card(slide, 6.76, 1.55, 2.54, 2.86, "实际价值", [
        "让分析顺序更清晰。",
        "让重点范围更收敛。",
        "让输出结果更容易被继续处理。",
    ], fill=SOFT_GREEN, accent=TEAL)

    # 23 Closing
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, 23, "总结", "项目汇报总结", "最后收回到一句话：平台的核心价值是辅助，而不是替代。")
    add_card(slide, 0.72, 1.6, 8.58, 2.95, fill=PANEL)
    add_textbox(slide, 0.98, 1.94, 7.9, 0.38, "本项目将序时账分析流程拆成“先理解数据、再校准规则、再形成样本、再辅助核实”的闭环。", size=18, color=NAVY, bold=True)
    add_multiline(slide, 1.0, 2.58, 7.85, 1.4, [
        "可视化部分负责建立整体判断和重点定位。",
        "规则校准部分负责把抽样参数调到更贴近当前账套特征。",
        "抽样与 LLM 辅助核实部分负责把结果收缩成更便于人工处理的样本集合。",
        "整体定位仍然是辅助分析工具，最终判断仍需回到审计证据与人工复核。",
    ], bullet=True, size=13.2)

    return prs


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
